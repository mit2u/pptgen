from pptx.util import Inches
from rest_framework import serializers
from pptx import Presentation
from core.models import SlideGenerationRequest
from core.utils import PPTX_MAPPING, SlideEnum
from pptx.dml.color import RGBColor

class SlideGenerationSerializer(serializers.ModelSerializer):
    ppt_link = serializers.SerializerMethodField()

    class Meta:
        model = SlideGenerationRequest
        fields = ("ppt_link","topic","content","font","color","min_slides","max_slides")



    def get_ppt_link(self, obj):
        slides_json = obj.slides

        prs = Presentation()
        for slide_json in slides_json:
            type = slide_json['slide_type']
            slide_layout = prs.slide_layouts[PPTX_MAPPING[slide_json['slide_type']]]
            slide = prs.slides.add_slide(slide_layout)
            slide_shapes = slide.shapes
            title = slide_shapes.title
            if title:
                self.add_text(slide_json['title'], title, obj)
            if type in [SlideEnum.title,SlideEnum.bullet]:
                subtitle = slide.placeholders[1]
                self.add_text(slide_json['content'], subtitle, obj)
            if type == SlideEnum.bullet:
                subtitle = slide.placeholders[1]
                text_frame = subtitle.text_frame
                for bullet in slide_json['bullet_points']:
                    p = text_frame.add_paragraph()
                    self.set_formatting(bullet,  p, obj)
                    p.level = 1
            elif type == SlideEnum.column:
                cols = 2
                rows = len(slide_json['col1_bullet_points'])
                left = top = Inches(2.0)
                width = Inches(6.0)
                height = Inches(0.8)

                table = slide_shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                table.columns[0].width = Inches(2.0)
                table.columns[1].width = Inches(4.0)

                self.add_column(obj, slide_json['col1_bullet_points'], table)
                self.add_column(obj, slide_json['col2_bullet_points'], table,1)
            elif type == SlideEnum.image:
                left = top = width = height = Inches(1)
                txBox = slide_shapes.add_textbox(left, top, width, height)
                self.add_text(slide_json['content'], txBox, obj)
            if slide_json['image']:
                img_path = slide_json['image']
                left = top = Inches(0.5)
                width = height = Inches(1)
                slide_shapes.add_picture(img_path, left, top,width, height)

        generate_ppt = str(obj.id) + ".pptx"
        prs.save(generate_ppt)
        request = self.context.get("request")
        return request.build_absolute_uri("/" + generate_ppt)

    def add_column(self, obj, slide_json, table,column=0):
        level = 0
        for bullet in slide_json:
            cell = table.cell(level, column)
            self.add_text(bullet, cell, obj)
            level += 1

    def set_formatting(self, bullet,  p, obj):
        run = p.add_run()
        run.text = bullet
        obj_font = "Arial"
        if obj.font:
            obj_font = obj.font
        rgb_color = RGBColor(0, 0, 0)
        if len(obj.color):
            try:
                rgb_color = RGBColor.from_string(obj.color)
            except:
                pass
        run.font.color.rgb = rgb_color
        run.font.name = obj_font

    def add_text(self, bullet, cell, obj):
        text_frame = cell.text_frame
        # Access the first paragraph within the text frame (assuming a single paragraph for simplicity)
        p = text_frame.paragraphs[0]
        # Access the run(s) within the paragraph to apply font formatting
        # If the cell already has text, you might need to iterate through existing runs,
        # or simply add a new run and set its text and font properties.
        self.set_formatting(bullet, p, obj)
